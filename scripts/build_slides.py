"""Build the CMPE 255 demo presentation deck (7 slides total).

Output: deliverables/slides/presentation.pptx

Designed for a timeboxed live demo (~10 minutes presenting + 5 minutes
demo + 5 minutes Q&A). Slide subset:

    0. Title
    1. Problem and what we're solving
    2. System at a glance (architecture)
    3. Headline results (table + alpha sweep)
    4. Live demo (screenshot anchor + presenter script)
    5. What didn't work (ablation, honesty)
    6. Conclusion + Q&A
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
RESULTS = ROOT / "data" / "results"
SCREENSHOTS = ROOT / "docs" / "screenshots"
OUT = ROOT / "deliverables" / "slides" / "presentation.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Palette
NAVY = RGBColor(0x0A, 0x2A, 0x4F)
NAVY_DEEP = RGBColor(0x06, 0x1A, 0x33)
ACCENT = RGBColor(0x3A, 0x86, 0xFF)
INK = RGBColor(0x1A, 0x23, 0x33)
MUTED = RGBColor(0x6B, 0x72, 0x80)
SOFT = RGBColor(0xE5, 0xEC, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLATE = RGBColor(0xCB, 0xD5, 0xE1)
GREEN = RGBColor(0x10, 0xB9, 0x81)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
ROSE = RGBColor(0xEF, 0x47, 0x6F)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)


def _para(tf, text, size=14, bold=False, color=INK, align=None,
          font="Calibri", first=False, space_after=4):
    para = tf.paragraphs[0] if first else tf.add_paragraph()
    if align is not None:
        para.alignment = align
    para.space_after = Pt(space_after)
    run = para.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return para


def add_text(slide, x, y, w, h, text, **kw):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    if isinstance(text, str):
        _para(tf, text, first=True, **kw)
    else:
        for i, t in enumerate(text):
            _para(tf, t, first=(i == 0), **kw)
    return box


def add_bullets(slide, x, y, w, h, bullets, size=14, color=INK,
                bullet_color=None):
    bc = bullet_color or ACCENT
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.space_after = Pt(7)
        # bullet glyph
        run0 = para.add_run()
        run0.text = "  "
        run0.font.size = Pt(size)
        run1 = para.add_run()
        run1.text = "■  "  # filled square
        run1.font.size = Pt(size - 2)
        run1.font.color.rgb = bc
        # body
        run2 = para.add_run()
        run2.text = b
        run2.font.name = "Calibri"
        run2.font.size = Pt(size)
        run2.font.color.rgb = color
    return box


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    return shape


def add_rounded(slide, x, y, w, h, fill=None, line=None, line_w=1.0):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = 0.08
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    return shape


def header(slide, eyebrow, title, subtitle=None):
    """Modern header: thin accent bar, eyebrow label, big title, subtitle."""
    bar = add_rect(slide, Inches(0.5), Inches(0.5), Inches(0.08),
                    Inches(0.55), fill=ACCENT)
    add_text(slide, Inches(0.7), Inches(0.5), Inches(8), Inches(0.3),
              eyebrow.upper(), size=10, bold=True, color=ACCENT)
    add_text(slide, Inches(0.7), Inches(0.78), Inches(12), Inches(0.55),
              title, size=28, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Inches(0.7), Inches(1.32), Inches(12), Inches(0.4),
                  subtitle, size=14, color=MUTED)
    # divider
    add_rect(slide, Inches(0.5), Inches(1.78), Inches(12.4), Inches(0.015),
              fill=SOFT)


def footer(slide, label):
    add_rect(slide, Inches(0.5), SLIDE_H - Inches(0.42), Inches(12.4),
              Inches(0.015), fill=SOFT)
    add_text(slide, Inches(0.5), SLIDE_H - Inches(0.36), Inches(7),
              Inches(0.3),
              "CMPE 255 Spring 2026  |  Patil, Kumar, Madhave",
              size=9, color=MUTED)
    add_text(slide, Inches(8.5), SLIDE_H - Inches(0.36), Inches(4.4),
              Inches(0.3), label, size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def picture_safe(slide, path: Path, x, y, **kw):
    if not path.exists():
        return None
    return slide.shapes.add_picture(str(path), x, y, **kw)


# --- Slide 0: Title -----------------------------------------------------

def slide_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # Full bleed navy gradient look (single deep navy + lighter band)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY_DEEP)
    add_rect(s, 0, Inches(5.0), SLIDE_W, Inches(2.5), fill=NAVY)
    # Accent stripe
    add_rect(s, 0, Inches(2.85), SLIDE_W, Inches(0.04), fill=ACCENT)

    # Eyebrow
    add_text(s, Inches(0.9), Inches(2.25), Inches(8), Inches(0.4),
              "CMPE 255 DATA MINING  |  SPRING 2026  |  FINAL PROJECT",
              size=12, bold=True, color=ACCENT)

    # Title
    add_text(s, Inches(0.9), Inches(3.10), Inches(11.5), Inches(1.2),
              "Record-Based Medical",
              size=52, bold=True, color=WHITE)
    add_text(s, Inches(0.9), Inches(3.95), Inches(11.5), Inches(1.2),
              "Diagnostic Assistant",
              size=52, bold=True, color=WHITE)

    # Tag line
    add_text(s, Inches(0.9), Inches(4.95), Inches(11.5), Inches(0.5),
              "FP-Growth association rules + biomedical RAG, fused.",
              size=18, color=SLATE)

    # Authors band
    add_text(s, Inches(0.9), Inches(6.05), Inches(11.5), Inches(0.4),
              "Sakshat Nandkumar Patil   ·   "
              "Vineet Kumar   ·   Aishwarya Madhave",
              size=16, bold=True, color=WHITE)
    add_text(s, Inches(0.9), Inches(6.45), Inches(11.5), Inches(0.4),
              "Department of Computer Engineering, San Jose State University",
              size=12, color=SLATE)


# --- Slide 1: Problem ---------------------------------------------------

def slide_problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "Slide 01  /  The Problem",
            "Why diagnostic decision support is hard",
            "Two failure modes in deployed clinical-AI tools today")

    # Left: failure modes as paired cards
    def fail_card(x, y, color, title, body):
        add_rounded(s, x, y, Inches(5.8), Inches(1.95),
                     fill=WHITE, line=color, line_w=1.5)
        # Color sidebar
        add_rect(s, x, y, Inches(0.10), Inches(1.95), fill=color)
        add_text(s, x + Inches(0.30), y + Inches(0.20),
                  Inches(5.4), Inches(0.4),
                  title, size=15, bold=True, color=color)
        add_text(s, x + Inches(0.30), y + Inches(0.65),
                  Inches(5.4), Inches(1.25),
                  body, size=12, color=INK)

    fail_card(Inches(0.5), Inches(2.05), ROSE,
               "Hand-built ontologies",
               "Expensive to maintain, brittle to new vocabulary, "
               "and they lock you into one domain expert's mental model. "
               "Updates lag the literature by years.")
    fail_card(Inches(0.5), Inches(4.10), AMBER,
               "Black-box LLMs",
               "Confident-sounding outputs without citations. "
               "Hallucination is silent. A clinician cannot defend "
               "the answer to a peer or to a patient.")

    # Right: our pitch as a stat panel
    panel = add_rounded(s, Inches(7.0), Inches(2.05),
                         Inches(5.9), Inches(4.0),
                         fill=NAVY)
    add_rect(s, Inches(7.0), Inches(2.05), Inches(0.10), Inches(4.0),
              fill=ACCENT)
    add_text(s, Inches(7.25), Inches(2.20), Inches(5.5), Inches(0.4),
              "OUR APPROACH", size=11, bold=True, color=ACCENT)
    add_text(s, Inches(7.25), Inches(2.55), Inches(5.5), Inches(0.6),
              "Auditable by construction", size=22, bold=True, color=WHITE)
    add_text(s, Inches(7.25), Inches(3.15), Inches(5.5), Inches(2.85),
              "Every prediction is grounded in two complementary signals:\n"
              "(1) an FP-Growth rule mined from real patient records, and\n"
              "(2) a biomedical passage retrieved from MedQuAD.\n\n"
              "The clinician sees both. The clinician decides.",
              size=13, color=SLATE)

    # Bottom strip: scope
    add_rounded(s, Inches(0.5), Inches(6.20), Inches(12.4), Inches(0.65),
                 fill=SOFT)
    add_text(s, Inches(0.7), Inches(6.30), Inches(12.0), Inches(0.5),
              "Scope: 41 disease classes  ·  131 symptom tokens  "
              "·  4,920 patient transactions  "
              "·  24,063 indexed passages",
              size=12, bold=True, color=NAVY)
    footer(s, "1 / 7")


# --- Slide 2: Architecture ----------------------------------------------

def slide_architecture(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "Slide 02  /  Architecture",
            "System at a glance",
            "Three-tier stack covering all six proposal steps")

    # Three horizontal lanes: Client / Service / Data, with the six
    # proposal steps grouped under Service.
    lane_x = Inches(0.5)
    lane_w = Inches(12.4)

    # Lane 1: Client
    add_rounded(s, lane_x, Inches(2.05), lane_w, Inches(0.85),
                 fill=SOFT, line=PURPLE, line_w=1.5)
    add_text(s, Inches(0.7), Inches(2.18), Inches(2.5), Inches(0.3),
              "WEB CLIENT", size=10, bold=True, color=PURPLE)
    add_text(s, Inches(0.7), Inches(2.45), Inches(12), Inches(0.4),
              "Next.js 14  ·  Symptom picker  ·  Live alpha slider  "
              "·  Pipeline timeline  ·  Insights dashboard",
              size=12, color=INK)

    # Arrow down
    add_text(s, Inches(6.4), Inches(2.92), Inches(0.5), Inches(0.25),
              "▼", size=14, color=MUTED, align=PP_ALIGN.CENTER)

    # Lane 2: Service (the six proposal steps as small chips)
    add_rounded(s, lane_x, Inches(3.18), lane_w, Inches(2.15),
                 fill=WHITE, line=ACCENT, line_w=1.5)
    add_text(s, Inches(0.7), Inches(3.30), Inches(6), Inches(0.3),
              "FASTAPI MICROSERVICE  (Python 3.11, MPS-aware)",
              size=10, bold=True, color=ACCENT)

    def step_chip(x_in, y_in, num, label, owner, color):
        x, y = Inches(x_in), Inches(y_in)
        w, h = Inches(2.0), Inches(0.80)
        add_rounded(s, x, y, w, h, fill=WHITE, line=color, line_w=1.2)
        add_rect(s, x, y, Inches(0.06), h, fill=color)
        add_text(s, x + Inches(0.15), y + Inches(0.07),
                  Inches(1.85), Inches(0.25),
                  f"STEP {num}", size=8, bold=True, color=color)
        add_text(s, x + Inches(0.15), y + Inches(0.30),
                  Inches(1.85), Inches(0.30),
                  label, size=11, bold=True, color=INK)
        add_text(s, x + Inches(0.15), y + Inches(0.55),
                  Inches(1.85), Inches(0.22),
                  owner, size=9, color=MUTED)

    base_x = 0.7
    chip_w = 2.05
    gap = 0.06
    y_top = 3.65
    y_bot = 4.50
    step_chip(base_x + 0 * (chip_w + gap), y_top, 1, "Data ETL",
               "Sakshat", MUTED)
    step_chip(base_x + 1 * (chip_w + gap), y_top, 2, "FP-Growth Mining",
               "Sakshat", ACCENT)
    step_chip(base_x + 2 * (chip_w + gap), y_top, 3, "Knowledge Base",
               "Vineet", GREEN)
    step_chip(base_x + 3 * (chip_w + gap), y_top, 4, "RAG Pipeline",
               "Vineet, Aishwarya", ROSE)
    step_chip(base_x + 4 * (chip_w + gap), y_top, 5, "Hybrid Fusion",
               "Aishwarya", AMBER)
    step_chip(base_x + 5 * (chip_w + gap), y_top, 6, "Evaluation",
               "Aishwarya", PURPLE)

    # Lane 2.5: a fusion equation strip
    add_rounded(s, Inches(0.7), Inches(y_bot), Inches(12.0),
                 Inches(0.65), fill=NAVY)
    add_text(s, Inches(0.9), Inches(y_bot + 0.10),
              Inches(11.6), Inches(0.5),
              "FusedScore(d) = a · RetrievalSim(d) + (1 - a) "
              "· MiningConf(d)        Default a = 0.3",
              size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
              font="Cambria Math")

    # Arrow down
    add_text(s, Inches(6.4), Inches(5.36), Inches(0.5), Inches(0.25),
              "▼", size=14, color=MUTED, align=PP_ALIGN.CENTER)

    # Lane 3: Data tier
    add_rounded(s, lane_x, Inches(5.62), lane_w, Inches(0.85),
                 fill=SOFT, line=GREEN, line_w=1.5)
    add_text(s, Inches(0.7), Inches(5.75), Inches(2.5), Inches(0.3),
              "DATA TIER", size=10, bold=True, color=GREEN)
    add_text(s, Inches(0.7), Inches(6.02), Inches(12), Inches(0.4),
              "23,839 rules  ·  24,063 passages  ·  "
              "FAISS (offline) or Pinecone Serverless 3072d (production)  "
              "·  Azure OpenAI embeddings + GPT-5.3 explainer",
              size=12, color=INK)
    footer(s, "2 / 7")


# --- Slide 3: Headline results ------------------------------------------

def slide_results(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "Slide 03  /  Results",
            "Headline metrics, 200 synthetic test cases",
            "Fusion beats mining alone. Retrieval alone is unusable.")

    # Top row: 4 big stat cards
    def stat(x_in, y_in, color, label, value, sub):
        x, y = Inches(x_in), Inches(y_in)
        w, h = Inches(2.95), Inches(1.55)
        add_rounded(s, x, y, w, h, fill=WHITE, line=color, line_w=1.5)
        add_rect(s, x, y, w, Inches(0.10), fill=color)
        add_text(s, x + Inches(0.20), y + Inches(0.20),
                  w - Inches(0.40), Inches(0.30),
                  label, size=10, bold=True, color=color)
        add_text(s, x + Inches(0.20), y + Inches(0.50),
                  w - Inches(0.40), Inches(0.65),
                  value, size=30, bold=True, color=NAVY)
        add_text(s, x + Inches(0.20), y + Inches(1.10),
                  w - Inches(0.40), Inches(0.40),
                  sub, size=10, color=MUTED)

    stat(0.50, 2.05, ACCENT, "RECALL @ 1",   "0.825", "fused, MiniLM, a = 0.3")
    stat(3.65, 2.05, GREEN,  "RECALL @ 10",  "0.890", "+2.0 over mining-only")
    stat(6.80, 2.05, AMBER,  "MRR",          "0.857", "+2.8 over mining-only")
    stat(9.95, 2.05, PURPLE, "LATENCY p95",  "20 ms", "M3 Pro, default config")

    # Mid row: ablation table on left, alpha sweep chart on right
    rows = [
        ["Variant",              "Mode",        "R@1",   "R@10",  "MRR"],
        ["mining-only",          "a = 0.0",     "0.790", "0.870", "0.829"],
        ["MiniLM",               "retrieval",   "0.130", "0.180", "0.151"],
        ["MiniLM",               "fused 0.3",   "0.825", "0.890", "0.857"],
        ["MiniLM + syn",         "fused 0.3",   "0.820", "0.895", "0.857"],
        ["PubMedBERT + syn",     "fused 0.3",   "0.815", "0.875", "0.843"],
    ]
    table_shape = s.shapes.add_table(len(rows), 5,
                                       Inches(0.5), Inches(3.85),
                                       Inches(7.2), Inches(2.55))
    tbl = table_shape.table
    widths = [2.2, 1.5, 1.0, 1.1, 1.1]
    for c, w in enumerate(widths):
        tbl.columns[c].width = Inches(w)

    for i, row in enumerate(rows):
        for j, v in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = v
            for para in cell.text_frame.paragraphs:
                para.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
                for run in para.runs:
                    run.font.name = "Calibri"
                    run.font.size = Pt(11)
                    run.font.bold = (i == 0) or (i == 3)  # header + winner
                    if i == 0:
                        run.font.color.rgb = WHITE
                    elif i == 3:
                        run.font.color.rgb = NAVY
                    else:
                        run.font.color.rgb = INK
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
            elif i == 3:
                cell.fill.solid()
                cell.fill.fore_color.rgb = SOFT

    # Alpha sweep chart on right
    picture_safe(s, RESULTS / "alpha_sweep.png",
                  Inches(7.95), Inches(3.85), width=Inches(4.95))
    add_text(s, Inches(7.95), Inches(6.30), Inches(4.95), Inches(0.3),
              "Alpha sweep: plateau on [0.1, 0.4]; collapse for a >= 0.7.",
              size=10, color=MUTED, align=PP_ALIGN.CENTER)

    # Bottom takeaway
    add_rounded(s, Inches(0.5), Inches(6.65), Inches(12.4),
                 Inches(0.50), fill=NAVY)
    add_text(s, Inches(0.7), Inches(6.71), Inches(12.0), Inches(0.4),
              "Takeaway: fusion lifts R@1 by +3.5 over the strong "
              "mining-only baseline. The signals are complementary.",
              size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    footer(s, "3 / 7")


# --- Slide 4: Live demo -------------------------------------------------

def slide_demo(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "Slide 04  /  Live Demo",
            "Cardiac-event preset, end-to-end",
            "Switch to the browser; if it fails, the screenshot is the demo")

    # Left: screenshot anchor
    pic = picture_safe(s, SCREENSHOTS / "02-results-cardiac.png",
                        Inches(0.5), Inches(2.05), height=Inches(4.55))
    if pic is not None:
        add_text(s, Inches(0.5), Inches(6.65), Inches(7.5), Inches(0.3),
                  "Production demo: Azure OpenAI 3072d + Pinecone + "
                  "GPT-5.3 explainer.",
                  size=10, color=MUTED)

    # Right: presenter script as a numbered list
    add_rounded(s, Inches(8.20), Inches(2.05),
                 Inches(4.7), Inches(4.95),
                 fill=NAVY)
    add_text(s, Inches(8.40), Inches(2.18), Inches(4.4), Inches(0.3),
              "DEMO SCRIPT  (~5 min)", size=10, bold=True, color=ACCENT)
    add_text(s, Inches(8.40), Inches(2.50), Inches(4.4), Inches(0.4),
              "Cardiac event walkthrough", size=15, bold=True, color=WHITE)

    steps = [
        ("1.", "Click the Cardiac event preset."),
        ("2.", "Heart Attack ranks #1 (fused 0.831)."),
        ("3.", "Read the 4-section clinical card."),
        ("4.", "Click a chip for symptom gloss + synonyms."),
        ("5.", "Drag the alpha slider, rankings re-fuse."),
        ("6.", "Filter to MedlinePlus only."),
        ("7.", "Toggle to OpenAI explainer."),
        ("8.", "Open Insights tab for live latency."),
    ]
    box = s.shapes.add_textbox(Inches(8.40), Inches(3.00),
                                 Inches(4.4), Inches(3.95))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (n, t) in enumerate(steps):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.space_after = Pt(5)
        run0 = para.add_run()
        run0.text = n + "  "
        run0.font.size = Pt(12)
        run0.font.bold = True
        run0.font.color.rgb = ACCENT
        run1 = para.add_run()
        run1.text = t
        run1.font.size = Pt(12)
        run1.font.color.rgb = SLATE
    footer(s, "4 / 7")


# --- Slide 5: What didn't work ------------------------------------------

def slide_lessons(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "Slide 05  /  Lessons",
            "What didn't work, and what we changed",
            "Three real bugs caught during development, all regression-tested")

    def lesson(x_in, y_in, num, color, title, body):
        x, y = Inches(x_in), Inches(y_in)
        w, h = Inches(4.0), Inches(4.55)
        add_rounded(s, x, y, w, h, fill=WHITE, line=color, line_w=1.5)
        # number bubble
        add_rect(s, x, y, w, Inches(0.65), fill=color)
        add_text(s, x + Inches(0.25), y + Inches(0.13),
                  Inches(3.5), Inches(0.4),
                  f"#{num}  {title}",
                  size=14, bold=True, color=WHITE)
        add_text(s, x + Inches(0.25), y + Inches(0.85),
                  w - Inches(0.5), h - Inches(1.0),
                  body, size=12, color=INK)

    lesson(0.5, 2.05, 1, ROSE,
            "min_support = 0.01",
            "Check-in 3 used 0.01. 21 of 41 diseases had no rules, so "
            "fusion was useless for those classes.\n\n"
            "Fix: dropped to 0.005. 100% disease coverage. The overlap "
            "weighting in MiningScorer.score neutralises the noisier "
            "rule list.")

    lesson(4.65, 2.05, 2, AMBER,
            "Passage-text attribution",
            "We matched a passage to a disease using its question AND "
            "answer body. Every \"What is Chest Pain?\" article got "
            "attributed to every cardiopulmonary disease.\n\n"
            "Fix: focus first, question second, never the answer body.")

    lesson(8.80, 2.05, 3, PURPLE,
            "hav matched have",
            "Disease keyword index had \"hav\" (Hepatitis A virus). "
            "Substring matching pulled it out of the word \"have\", "
            "which is in basically every passage.\n\n"
            "Fix: word-boundary regex for single-token keywords. "
            "Multi-word phrases stay substring. Regression test in "
            "test_disease_keywords.py.")

    # Bottom: honesty strip
    add_rounded(s, Inches(0.5), Inches(6.75), Inches(12.4),
                 Inches(0.40), fill=SOFT)
    add_text(s, Inches(0.7), Inches(6.80), Inches(12.0), Inches(0.3),
              "Also tried and dropped: always-on cross-encoder "
              "(+100 ms for marginal gain), Cohere reranker by default "
              "(403 on free Pinecone tier).",
              size=10, color=MUTED, align=PP_ALIGN.CENTER)
    footer(s, "5 / 7")


# --- Slide 6: Conclusion + Q&A -----------------------------------------

def slide_conclusion(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "Slide 06  /  Conclusion",
            "What we shipped, and where it goes next",
            "")

    # Left column: shipped (with green check accent)
    add_text(s, Inches(0.5), Inches(2.05), Inches(6), Inches(0.4),
              "SHIPPED", size=12, bold=True, color=GREEN)
    shipped = [
        "Reproducible FP-Growth pipeline, 23,839 rules, 100% disease coverage",
        "Three encoder backends: MiniLM 384d, PubMedBERT 768d, Azure 3072d",
        "Pinecone Serverless production index (24,063 vectors, 3072d)",
        "Claim-level evidence cards with NIH source-tier badges",
        "Dual explainer: deterministic template + structured OpenAI",
        "155 pytest tests passing (147 unit, 8 live integration)",
        "82.5% R@1, 0.857 MRR. 12 ms p50 end-to-end on the offline path.",
    ]
    add_bullets(s, Inches(0.5), Inches(2.45), Inches(6.0), Inches(4.0),
                 shipped, size=12, bullet_color=GREEN)

    # Right column: future work
    add_text(s, Inches(6.85), Inches(2.05), Inches(6), Inches(0.4),
              "FUTURE WORK", size=12, bold=True, color=AMBER)
    future = [
        "UMLS-backed automatic synonym mining (replace curated dict)",
        "Cross-encoder over <disease, passage> pairs in the default path",
        "Held-out vocabulary eval on the Synthea FHIR pipeline",
        "Open-domain disease universe via UMLS canonicalisation",
        "Learned fusion (XGBoost) using rule lift, source tier, overlap",
    ]
    add_bullets(s, Inches(6.85), Inches(2.45), Inches(6.0), Inches(4.0),
                 future, size=12, bullet_color=AMBER)

    # Q&A banner
    add_rounded(s, Inches(0.5), Inches(6.30), Inches(12.4),
                 Inches(0.85), fill=NAVY)
    add_text(s, Inches(0.5), Inches(6.42), Inches(12.4), Inches(0.4),
              "Questions?", size=22, bold=True, color=WHITE,
              align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.85), Inches(12.4), Inches(0.3),
              "github.com/sakshat-patil/Symptom-Based-Disease-Identification-via-RAG",
              size=11, color=SLATE, align=PP_ALIGN.CENTER)
    footer(s, "6 / 7")


# --- Build --------------------------------------------------------------

def build():
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)         # cover (not numbered as 0/7)
    slide_problem(prs)       # 1 / 7
    slide_architecture(prs)  # 2 / 7
    slide_results(prs)       # 3 / 7
    slide_demo(prs)          # 4 / 7
    slide_lessons(prs)       # 5 / 7
    slide_conclusion(prs)    # 6 / 7

    prs.save(str(OUT))
    print(f"[slides] wrote {OUT}")


if __name__ == "__main__":
    build()
